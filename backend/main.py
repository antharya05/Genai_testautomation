import os
import shutil

import fitz
from docx import Document
from fastapi import FastAPI, UploadFile, File
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
from pptx import Presentation

import llm
from parser import parse_requirements

app = FastAPI(title="AI Test Case Generator")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

UPLOAD_FOLDER = "uploads"
os.makedirs(UPLOAD_FOLDER, exist_ok=True)

ALLOWED_EXTENSIONS = (".pdf", ".docx", ".pptx", ".txt")


# --------------------------------------------------
# Pydantic models
# --------------------------------------------------

class GenerateRequest(BaseModel):
    requirements: list[str]


class TextRequest(BaseModel):
    text: str


class TestCase(BaseModel):
    id: str
    requirement_id: str
    title: str
    test_type: str
    asil_level: str
    preconditions: list[str]
    steps: list[str]
    expected_result: str
    priority: str


# --------------------------------------------------
# Text extractors
# --------------------------------------------------

def extract_pdf(path: str) -> str:
    doc = fitz.open(path)
    return "".join(page.get_text() for page in doc)


def extract_docx(path: str) -> str:
    doc = Document(path)
    return "\n".join(p.text for p in doc.paragraphs)


def extract_pptx(path: str) -> str:
    ppt = Presentation(path)
    parts = []
    for slide in ppt.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                parts.append(shape.text)
    return "\n".join(parts)


def extract_txt(path: str) -> str:
    with open(path, "r", encoding="utf-8") as f:
        return f.read()


EXTRACTOR_MAP = {
    ".pdf": extract_pdf,
    ".docx": extract_docx,
    ".pptx": extract_pptx,
    ".txt": extract_txt,
}


# --------------------------------------------------
# Routes
# --------------------------------------------------

@app.get("/health")
def health():
    return {"status": "ok"}


@app.post("/upload")
async def upload_file(file: UploadFile = File(...)):
    ext = os.path.splitext(file.filename.lower())[1]

    if ext not in ALLOWED_EXTENSIONS:
        return {"error": f"Unsupported file type '{ext}'. Allowed: {', '.join(ALLOWED_EXTENSIONS)}"}

    file_path = os.path.join(UPLOAD_FOLDER, file.filename)
    with open(file_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)

    extractor = EXTRACTOR_MAP[ext]
    extracted_text = extractor(file_path)
    requirements = parse_requirements(extracted_text)

    return {
        "filename": file.filename,
        "extracted_text": extracted_text,
        "requirements": requirements,
        "requirement_count": len(requirements),
    }


@app.post("/parse-text")
async def parse_text(request: TextRequest):
    if not request.text or not request.text.strip():
        return {"error": "No text provided", "requirements": []}
    requirements = parse_requirements(request.text)
    return {
        "filename": "Pasted Text",
        "extracted_text": request.text,
        "requirements": requirements,
        "requirement_count": len(requirements),
    }


@app.post("/generate")
async def generate(request: GenerateRequest):
    if not request.requirements:
        return {"error": "No requirements provided", "test_cases": []}

    try:
        test_cases = llm.generate_test_cases(request.requirements)
    except Exception as e:
        return {"error": f"AI generation failed: {str(e)}", "test_cases": []}

    return {
        "test_cases": test_cases,
        "count": len(test_cases),
    }
